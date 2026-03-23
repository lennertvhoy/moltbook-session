"""Render low-fidelity PNG previews from the generated Moltbook.pptx.

The environment does not ship with PowerPoint or LibreOffice, so this script
uses python-pptx plus Pillow to render layout-faithful previews directly from
the generated PPTX. The previews are sufficient for slide-by-slide QA on text
fit, spacing, image placement, and visual hierarchy.
"""

from __future__ import annotations

import io
import json
from pathlib import Path
from textwrap import wrap
from zipfile import ZipFile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "Moltbook.pptx"
OUT_DIR = ROOT / "qa" / "previews"
CONTACT_SHEET = OUT_DIR / "contact-sheet.png"
META_PATH = OUT_DIR / "preview-metadata.json"

BACKGROUND = "#F6F1E8"
FONT_PATH = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
FONT_BOLD_PATH = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
PX_WIDTH = 1600
PX_HEIGHT = 900


def emu_to_px(value: int, total_emu: int, total_px: int) -> int:
    """Convert EMU coordinates into pixels."""
    return round(value / total_emu * total_px)


def load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    """Load a truetype font if available, else fall back."""
    try:
        path = FONT_BOLD_PATH if bold else FONT_PATH
        return ImageFont.truetype(path, size=size)
    except OSError:
        return ImageFont.load_default()


def draw_text_box(
    draw: ImageDraw.ImageDraw,
    text: str,
    box: tuple[int, int, int, int],
    *,
    font_size: int = 24,
    color: str = "#17212B",
    bold: bool = False,
) -> None:
    """Draw wrapped text into the given box."""
    if not text.strip():
        return

    x, y, w, h = box
    font = load_font(font_size, bold=bold)
    avg_char_width = max(font_size * 0.52, 6)
    max_chars = max(int(w / avg_char_width), 1)
    wrapped_lines: list[str] = []
    for raw_line in text.splitlines():
        if not raw_line.strip():
            wrapped_lines.append("")
            continue
        wrapped_lines.extend(wrap(raw_line, width=max_chars, replace_whitespace=False))

    line_height = int(font_size * 1.3)
    max_lines = max(int(h / line_height), 1)
    visible_lines = wrapped_lines[:max_lines]
    cursor_y = y
    for line in visible_lines:
        draw.text((x, cursor_y), line, font=font, fill=color)
        cursor_y += line_height


def shape_fill_color(shape) -> str | None:
    """Best-effort shape fill color extraction."""
    fill = shape.fill
    if fill is None or fill.type is None:
        return None
    try:
        fore_color = getattr(fill, "fore_color", None)
    except TypeError:
        return None
    if fore_color is None:
        return None
    rgb = getattr(fore_color, "rgb", None)
    return f"#{rgb}" if rgb else None


def shape_line_color(shape) -> str:
    """Best-effort shape line color extraction."""
    line = getattr(shape, "line", None)
    if line is None:
        return "#000000"
    color = getattr(line, "color", None)
    if color is None:
        fore = getattr(line, "fore_color", None)
        rgb = getattr(fore, "rgb", None) if fore else None
        return f"#{rgb}" if rgb else "#000000"
    rgb = getattr(color, "rgb", None)
    return f"#{rgb}" if rgb else "#000000"


def render_slide(slide, slide_width: int, slide_height: int, index: int) -> dict[str, object]:
    """Render one slide to a PNG preview."""
    canvas = Image.new("RGB", (PX_WIDTH, PX_HEIGHT), BACKGROUND)
    draw = ImageDraw.Draw(canvas)

    title = ""
    text_runs = 0
    image_count = 0
    top_texts: list[str] = []

    for shape in slide.shapes:
        left = emu_to_px(shape.left, slide_width, PX_WIDTH)
        top = emu_to_px(shape.top, slide_height, PX_HEIGHT)
        width = emu_to_px(shape.width, slide_width, PX_WIDTH)
        height = emu_to_px(shape.height, slide_height, PX_HEIGHT)

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
            image = image.resize((max(width, 1), max(height, 1)))
            canvas.paste(image, (left, top))
            image_count += 1
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            fill_color = shape_fill_color(shape)
            line_color = shape_line_color(shape)
            if fill_color:
                radius = 18 if "ROUND" in str(shape.auto_shape_type) else 0
                if radius:
                    draw.rounded_rectangle(
                        (left, top, left + width, top + height),
                        radius=radius,
                        fill=fill_color,
                        outline=line_color,
                        width=2,
                    )
                else:
                    draw.rectangle(
                        (left, top, left + width, top + height),
                        fill=fill_color,
                        outline=line_color,
                        width=2,
                    )
            elif width > 0 and height <= 4:
                draw.line((left, top, left + width, top), fill=line_color, width=2)

        if hasattr(shape, "text_frame"):
            text = shape.text.strip()
            if not text:
                continue
            text_runs += len(text.split())
            if top < 140:
                top_texts.append(text)
            if not title and top < 150 and len(text) > 8:
                title = text

            is_kicker = top < 90 and len(text) < 40
            is_title = top < 150 and len(text) > 20
            font_size = 13 if is_kicker else 26 if is_title else 18
            color = "#0F766E" if is_kicker else "#17212B"
            draw_text_box(
                draw,
                text,
                (left, top, width, height),
                font_size=font_size,
                color=color,
                bold=is_kicker or is_title,
            )

    slide_name = f"slide-{index:02d}.png"
    output_path = OUT_DIR / slide_name
    canvas.save(output_path)

    return {
        "index": index,
        "title": title or f"Slide {index}",
        "preview": str(output_path.relative_to(ROOT)),
        "image_count": image_count,
        "text_word_count": text_runs,
        "top_texts": top_texts,
    }


def build_contact_sheet(preview_paths: list[Path]) -> None:
    """Assemble a contact sheet from individual slide previews."""
    thumbs = [Image.open(path).resize((480, 270)) for path in preview_paths]
    cols = 2
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * 500, rows * 300), "#EEE7DA")
    draw = ImageDraw.Draw(sheet)
    label_font = load_font(18, bold=True)

    for idx, thumb in enumerate(thumbs):
        row = idx // cols
        col = idx % cols
        x = col * 500 + 10
        y = row * 300 + 10
        sheet.paste(thumb, (x, y))
        draw.text((x, y + 276), f"Slide {idx + 1}", fill="#17212B", font=label_font)

    sheet.save(CONTACT_SHEET)


def count_notes() -> int:
    """Count notes slides in the generated PPTX."""
    with ZipFile(PPTX_PATH) as archive:
        return len([name for name in archive.namelist() if name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml")])


def main() -> None:
    """Render all slides and write preview metadata."""
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    presentation = Presentation(PPTX_PATH)
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    metadata = []
    preview_paths: list[Path] = []
    for index, slide in enumerate(presentation.slides, start=1):
        item = render_slide(slide, slide_width, slide_height, index)
        metadata.append(item)
        preview_paths.append(ROOT / str(item["preview"]))

    build_contact_sheet(preview_paths)
    payload = {
        "slide_count": len(metadata),
        "notes_count": count_notes(),
        "slides": metadata,
        "contact_sheet": str(CONTACT_SHEET.relative_to(ROOT)),
    }
    META_PATH.write_text(json.dumps(payload, indent=2), encoding="utf-8")
    print(json.dumps(payload, indent=2))


if __name__ == "__main__":
    main()
